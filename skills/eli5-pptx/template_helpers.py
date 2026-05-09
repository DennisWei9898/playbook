"""
ELI5 PPTX Template Helpers — 可重用的 python-pptx 工具函式庫
配合 playbook/skills/eli5-pptx/SKILL.md 使用

使用方式：
    import sys
    sys.path.insert(0, '/path/to/playbook/skills/eli5-pptx')
    from template_helpers import *

    prs, P = init_presentation('polymarket')  # P = palette dict
    # 用 P['ACCENT'] / P['BG_CREAM'] 等取色
    # 用 add_text(slide, ..., color=P['TEXT_DARK']) 等
    prs.save('output.pptx')
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ════════════════════════════════════════════════
# 4 個預設配色 Palette
# ════════════════════════════════════════════════

def _hex(s):
    """#faf6f1 → RGBColor(0xfa, 0xf6, 0xf1)"""
    s = s.lstrip('#')
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


PALETTE_POLYMARKET = {
    'name':         'polymarket',
    'BG_CREAM':     _hex('#faf6f1'),
    'BG_WARM':      _hex('#f5ede2'),
    'BG_DARK':      _hex('#2a1208'),
    'BG_BROWN':     _hex('#3a1c10'),
    'TEXT_DARK':    _hex('#2a1208'),
    'TEXT_BROWN':   _hex('#7a4a38'),
    'TEXT_TAUPE':   _hex('#b09080'),
    'ACCENT':       _hex('#c85a3c'),
    'PEACH':        _hex('#f7d5c8'),
    'PEACH_DEEP':   _hex('#eda890'),
    'GREEN':        _hex('#3a9e8f'),
    'GREEN_LIGHT':  _hex('#d4f0ec'),
    'AMBER':        _hex('#b07800'),
    'AMBER_LIGHT':  _hex('#f5c66b'),
    'DANGER':       _hex('#b83a2c'),
    'DANGER_LIGHT': _hex('#f6c8c0'),
    'PURPLE':       _hex('#8b5cf6'),
    'WHITE':        _hex('#ffffff'),
    'BORDER':       _hex('#e5d5c4'),
}

PALETTE_MCKINSEY = {
    'name':         'mckinsey',
    'BG_CREAM':     _hex('#ffffff'),
    'BG_WARM':      _hex('#f4f6f8'),
    'BG_DARK':      _hex('#0f2c4a'),
    'BG_BROWN':     _hex('#1a3a5e'),
    'TEXT_DARK':    _hex('#0f2c4a'),
    'TEXT_BROWN':   _hex('#4a5a6a'),
    'TEXT_TAUPE':   _hex('#a0aab4'),
    'ACCENT':       _hex('#1f6fb4'),
    'PEACH':        _hex('#c8dff0'),
    'PEACH_DEEP':   _hex('#7fb3df'),
    'GREEN':        _hex('#2ca58d'),
    'GREEN_LIGHT':  _hex('#cce8e0'),
    'AMBER':        _hex('#d09000'),
    'AMBER_LIGHT':  _hex('#f5d870'),
    'DANGER':       _hex('#c0392b'),
    'DANGER_LIGHT': _hex('#f4c8c0'),
    'PURPLE':       _hex('#7e57c2'),
    'WHITE':        _hex('#ffffff'),
    'BORDER':       _hex('#dde3e8'),
}

PALETTE_NOTION = {
    'name':         'notion',
    'BG_CREAM':     _hex('#ffffff'),
    'BG_WARM':      _hex('#f7f7f5'),
    'BG_DARK':      _hex('#191919'),
    'BG_BROWN':     _hex('#2f2f2f'),
    'TEXT_DARK':    _hex('#191919'),
    'TEXT_BROWN':   _hex('#4f4f4f'),
    'TEXT_TAUPE':   _hex('#9b9b9b'),
    'ACCENT':       _hex('#000000'),
    'PEACH':        _hex('#e8e8e6'),
    'PEACH_DEEP':   _hex('#c8c8c4'),
    'GREEN':        _hex('#0f7b6c'),
    'GREEN_LIGHT':  _hex('#d4ebe4'),
    'AMBER':        _hex('#b07800'),
    'AMBER_LIGHT':  _hex('#f5d870'),
    'DANGER':       _hex('#cc4646'),
    'DANGER_LIGHT': _hex('#f4d8d8'),
    'PURPLE':       _hex('#6940a5'),
    'WHITE':        _hex('#ffffff'),
    'BORDER':       _hex('#e8e8e6'),
}

PALETTES = {
    'polymarket': PALETTE_POLYMARKET,
    'mckinsey':   PALETTE_MCKINSEY,
    'notion':     PALETTE_NOTION,
}


def get_palette(name='polymarket', custom=None):
    """
    取得配色。name 三選一，或傳 custom dict 自訂（必須含所有 key）。
    """
    if custom is not None:
        # 自訂配色：missing keys fall back to polymarket
        merged = dict(PALETTE_POLYMARKET)
        for k, v in custom.items():
            if isinstance(v, str):
                merged[k] = _hex(v)
            else:
                merged[k] = v
        merged['name'] = 'custom'
        return merged
    return PALETTES.get(name, PALETTE_POLYMARKET)


# ════════════════════════════════════════════════
# Presentation 初始化
# ════════════════════════════════════════════════

def init_presentation(palette_name='polymarket', custom_palette=None):
    """
    建立 16:9 Presentation 並回傳 (prs, palette dict)。
    用法：
        prs, P = init_presentation('polymarket')
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        set_bg(slide, P['BG_CREAM'])
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    palette = get_palette(palette_name, custom_palette)
    return prs, palette


# ════════════════════════════════════════════════
# 工具函數
# ════════════════════════════════════════════════

def set_bg(slide, color):
    """設背景色（接受 RGBColor）"""
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Noto Sans TC", italic=False):
    """新增單行（或單段）文字。color 預設為黑（呼叫者通常會傳 P['TEXT_DARK']）。"""
    if color is None:
        color = RGBColor(0, 0, 0)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_multiline(slide, x, y, w, h, lines, *, default_size=12, default_color=None,
                  line_spacing=1.4, anchor=MSO_ANCHOR.TOP):
    """
    多行文字，每行可單獨 style。
    lines: list of (text, opts) 或 list of str
    opts 可含 size / bold / italic / color / font / align / line_spacing
    """
    if default_color is None:
        default_color = RGBColor(0, 0, 0)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        text, opts = (item if isinstance(item, tuple) else (item, {}))
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get('align', PP_ALIGN.LEFT)
        p.line_spacing = opts.get('line_spacing', line_spacing)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get('size', default_size))
        run.font.bold = opts.get('bold', False)
        run.font.italic = opts.get('italic', False)
        run.font.color.rgb = opts.get('color', default_color)
        run.font.name = opts.get('font', 'Noto Sans TC')
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75):
    """方形矩形。"""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.text_frame.margin_left = Emu(0)
    sh.text_frame.margin_right = Emu(0)
    sh.text_frame.margin_top = Emu(0)
    sh.text_frame.margin_bottom = Emu(0)
    return sh


def add_rounded(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75, corner=0.05):
    """圓角矩形。corner 0~0.5 控制圓角強度。"""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    sh.adjustments[0] = corner
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.text_frame.margin_left = Emu(0)
    sh.text_frame.margin_right = Emu(0)
    sh.text_frame.margin_top = Emu(0)
    sh.text_frame.margin_bottom = Emu(0)
    return sh


def add_oval(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75):
    """橢圓 / 圓。"""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    return sh


# ════════════════════════════════════════════════
# 標準頁面元件
# ════════════════════════════════════════════════

class PageCounter:
    """
    讓 page_footer 可以自動編頁。
    用法：
        page_counter = PageCounter(total=16)
        page_footer(slide, "PART 1  ·  Story", palette, page_counter)
    """
    def __init__(self, total):
        self.current = 0
        self.total = total

    def next(self):
        self.current += 1
        return self.current


def page_footer(slide, label, palette, page_counter):
    """
    左下角 PART label（accent 色），右下頁碼（taupe 色）。
    """
    page_counter.next()
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             label, size=10, bold=True, color=palette['ACCENT'])
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{page_counter.current} / {page_counter.total}",
             size=10, bold=True, color=palette['TEXT_TAUPE'],
             align=PP_ALIGN.RIGHT)


def section_title(slide, label_top, title_text, palette, subtitle=""):
    """
    每頁頂部標準標題區：
    - PART X · Y 小標籤（accent 色）
    - 大標題（深色）
    - 副標（褐色 / 次要色）
    """
    add_text(slide, Inches(0.6), Inches(0.55), Inches(11.5), Inches(0.3),
             label_top, size=11, bold=True, color=palette['ACCENT'])
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12.0), Inches(0.7),
             title_text, size=32, bold=True, color=palette['TEXT_DARK'])
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.55), Inches(12.0), Inches(0.4),
                 subtitle, size=15, color=palette['TEXT_BROWN'])


# ════════════════════════════════════════════════
# 常用組件（高層級 helper）
# ════════════════════════════════════════════════

def add_blank_slide(prs, palette):
    """加一張空白 slide 並設好背景。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, palette['BG_CREAM'])
    return slide


def add_metric_card(slide, x, y, w, h, label, big_text, sub_text, palette, *, accent_color=None):
    """
    指標卡：頂部小 label + 大數字 + 底部說明。
    用於「技術成本」「市場規模」這類頁面。
    """
    color = accent_color or palette['ACCENT']
    add_rounded(slide, x, y, w, h,
                fill=palette['BG_WARM'], line=color, line_w=2.0, corner=0.04)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.4),
             label, size=14, bold=True, color=color)
    if isinstance(big_text, str):
        big_lines = big_text.split("\n")
    else:
        big_lines = list(big_text)
    add_multiline(slide, x + Inches(0.3), y + Inches(0.85), w - Inches(0.6),
                  Inches(h.inches * 0.45),
                  [(line, {'size': 38, 'bold': True, 'color': palette['TEXT_DARK'],
                           'line_spacing': 1.1})
                   for line in big_lines])
    if isinstance(sub_text, str):
        sub_lines = sub_text.split("\n")
    else:
        sub_lines = list(sub_text)
    add_multiline(slide, x + Inches(0.3), y + Inches(h.inches * 0.6),
                  w - Inches(0.6), Inches(h.inches * 0.4),
                  [(line, {'size': 12, 'color': palette['TEXT_BROWN'], 'line_spacing': 1.5})
                   for line in sub_lines])


def add_quadrant_card(slide, x, y, w, h, emoji, title, body, palette, *, accent_color=None):
    """
    四象限卡片：emoji + 標題 + 內文。
    用於「Problem Statement Canvas」「競品全景」這類 4 格頁面。
    """
    color = accent_color or palette['ACCENT']
    add_rounded(slide, x, y, w, h,
                fill=palette['BG_WARM'], line=color, line_w=2.0, corner=0.04)
    add_text(slide, x + Inches(0.2), y + Inches(0.25), w - Inches(0.4), Inches(0.7),
             emoji, size=42)
    add_text(slide, x + Inches(0.2), y + Inches(1.15), w - Inches(0.4), Inches(0.4),
             title, size=18, bold=True, color=color)
    if isinstance(body, str):
        body_lines = body.split("\n")
    else:
        body_lines = list(body)
    add_multiline(slide, x + Inches(0.2), y + Inches(1.7), w - Inches(0.4),
                  Inches(h.inches - 1.9),
                  [(line, {'size': 13, 'color': palette['TEXT_DARK'], 'line_spacing': 1.5})
                   for line in body_lines])


def add_pricing_card(slide, x, y, w, h, name, price, target, features, palette,
                     *, accent_color=None, highlight=False):
    """
    定價方案卡片：名稱 + 價格 + 對象 + 功能列表。
    用於「商業模式」這頁。
    """
    color = accent_color or palette['ACCENT']
    line_w = 3.0 if highlight else 1.0
    bg = palette['PEACH'] if highlight else palette['BG_WARM']
    add_rounded(slide, x, y, w, h, fill=bg, line=color, line_w=line_w, corner=0.04)
    # 方案名
    add_text(slide, x + Inches(0.2), y + Inches(0.25), w - Inches(0.4), Inches(0.5),
             name, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    # 價格
    if isinstance(price, str):
        price_lines = price.split("\n")
    else:
        price_lines = list(price)
    add_multiline(slide, x + Inches(0.2), y + Inches(0.85), w - Inches(0.4), Inches(1.0),
                  [(line, {'size': 18, 'bold': True, 'color': palette['TEXT_DARK'],
                           'align': PP_ALIGN.CENTER, 'line_spacing': 1.3})
                   for line in price_lines])
    # 對象
    add_text(slide, x + Inches(0.2), y + Inches(2.05), w - Inches(0.4), Inches(0.4),
             target, size=11, italic=True, color=palette['TEXT_BROWN'], align=PP_ALIGN.CENTER)
    # 分隔線
    add_rect(slide, x + Inches(0.4), y + Inches(2.55), w - Inches(0.8), Inches(0.02), fill=color)
    # 功能列表
    if isinstance(features, str):
        feature_lines = features.split("\n")
    else:
        feature_lines = list(features)
    add_multiline(slide, x + Inches(0.25), y + Inches(2.7), w - Inches(0.5),
                  Inches(h.inches - 2.9),
                  [(line, {'size': 11, 'color': palette['TEXT_DARK'], 'line_spacing': 1.4})
                   for line in feature_lines])


# ════════════════════════════════════════════════
# Re-export 常用 import
# ════════════════════════════════════════════════
__all__ = [
    # python-pptx 直通
    'Presentation', 'Inches', 'Pt', 'Emu', 'RGBColor',
    'MSO_SHAPE', 'PP_ALIGN', 'MSO_ANCHOR',
    # palettes
    'PALETTE_POLYMARKET', 'PALETTE_MCKINSEY', 'PALETTE_NOTION', 'PALETTES',
    'get_palette',
    # init
    'init_presentation', 'add_blank_slide',
    # primitives
    'set_bg', 'add_text', 'add_multiline',
    'add_rect', 'add_rounded', 'add_oval',
    # standard components
    'PageCounter', 'page_footer', 'section_title',
    # high-level cards
    'add_metric_card', 'add_quadrant_card', 'add_pricing_card',
]
